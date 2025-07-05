from pptx.enum.shapes import MSO_CONNECTOR
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor


# Create a new presentation
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Workflow: Staff Scheduling → Payroll System"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.bold = True

# Step 1: Staff Scheduling System,
box1_left = Inches(0.5)
box1_top = Inches(1.5)
box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box1_left, box1_top, Inches(3), Inches(1))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0, 102, 204)
box1.text_frame.text = "Use Case 1:\nStaff Scheduling System"
for p in box1.text_frame.paragraphs:
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True

# Step 2: Integration Layer
box2_left = Inches(4.5)
box2_top = Inches(1.5)
box2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box2_left, box2_top, Inches(3), Inches(1))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(255, 204, 0)
box2.text_frame.text = "Integration Layer:\nExport shift data (e.g., hours, overtime)"
for p in box2.text_frame.paragraphs:
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.bold = True

# Step 3: Payroll System..
box3_left = Inches(8.5)
box3_top = Inches(1.5)
box3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box3_left, box3_top, Inches(3), Inches(1))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0, 153, 76)
box3.text_frame.text = "Use Case 2:\nPayroll/Salary Management System"
for p in box3.text_frame.paragraphs:
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True

# Connectors
conn1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, box1_left + Inches(3), box1_top + Inches(0.5),
                                   box2_left, box2_top + Inches(0.5))
conn1.line.color.rgb = RGBColor(100, 100, 100)

conn2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, box2_left + Inches(3), box2_top + Inches(0.5),
                                   box3_left, box3_top + Inches(0.5))
conn2.line.color.rgb = RGBColor(100, 100, 100)

# Save the presentation
pptx_path = "/Users/bhavyavadher/Desktop/Staff_Scheduling_to_Payroll_Flow.pptx"
prs.save(pptx_path)

pptx_path